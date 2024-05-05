import sqlite3
import openai
import tkinter as tk
from tkinter import messagebox, filedialog
from docx import Document
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation

# Read the API key from the configuration file
CONFIG_FILE_PATH = 'config.txt'

def read_api_key():
    """
    Read the API key from the config.txt file.

    Returns:
    - str: The API key read from the file.
    """
    try:
        with open(CONFIG_FILE_PATH, 'r') as file:
            return file.readline().strip()
    except FileNotFoundError:
        messagebox.showerror("Configuration Error", "The config.txt file was not found.")
        return None

# Set the OpenAI API key
openai.api_key = read_api_key()

class ChatGPTApp:
    def __init__(self, root):
        self.root = root
        self.root.title("ChatGPT")

        # Initialize the database connection
        self.conn = sqlite3.connect('chatgpt_responses.db')
        self.create_table_if_not_exists()

        # Initialize class attribute for uploaded file data
        self.uploaded_file_data = None

        # Initialize the GUI
        self.setup_gui()

    def setup_gui(self):
        # Create a label for the user input prompt
        prompt_label = tk.Label(self.root, text="Enter your prompt for ChatGPT:")
        prompt_label.pack(pady=10)

        # Create an entry widget for the user's prompt
        self.entry = tk.Entry(self.root, width=50)
        self.entry.pack(pady=5)
        self.entry.bind("<Return>", self.on_submit)  # Trigger on submission (Enter key press)

        # Create a text widget for displaying ChatGPT responses and file data
        self.text_output = tk.Text(self.root, width=60, height=20)
        self.text_output.pack(pady=10)

        # Create a button for submitting the user's prompt
        submit_button = tk.Button(self.root, text="Submit", command=self.on_submit)
        submit_button.pack(pady=10)

        # Create a button for viewing previous conversations
        view_button = tk.Button(self.root, text="View Previous Conversations", command=self.view_previous_conversations)
        view_button.pack(pady=10)

        # Create a button for uploading files
        upload_button = tk.Button(self.root, text="Upload File", command=self.upload_file)
        upload_button.pack(pady=10)

    def on_submit(self, event=None):
        """
        Handle user input and generate ChatGPT response.
        """
        user_input = self.entry.get()
        if user_input.lower() == 'exit':
            self.close_app()
            return
        
        # Include uploaded file data in the prompt if available
        prompt = user_input
        if self.uploaded_file_data:
            prompt += f"\n\nFile data:\n{self.uploaded_file_data}"

        try:
            # Generate ChatGPT response using the prompt
            response = self.generate_response(prompt)
            success = self.insert_conversation(user_input, response, "")

            # Handle potential database insertion errors
            if not success:
                messagebox.showerror("Database Error", "Failed to insert conversation into the database.")
            
            # Display user input and ChatGPT response in the text widget
            self.text_output.insert(tk.END, f"You: {user_input}\n")
            self.text_output.insert(tk.END, f"ChatGPT: {response}\n")
            self.entry.delete(0, tk.END)
        except Exception as e:
            messagebox.showerror("Error", str(e))

    def generate_response(self, prompt):
        """
        Generate a response from the OpenAI API given a prompt.
        """
        try:
            # Call the ChatGPT API to generate a response
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=150
            )
            # Extract the text response from the API result
            return response['choices'][0]['message']['content'].strip()
        except Exception as e:
            # Handle API errors and display an error message
            messagebox.showerror("API Error", f"Failed to connect to OpenAI API: {str(e)}")
            return "Unable to retrieve response from ChatGPT at the moment."

    def upload_file(self):
        """
        Handle file uploads and process file contents.
        """
        try:
            # Add '.pdf' and '.ppt' filetypes to the file dialog
            filetypes=[("All Files", "*.*"), 
                       ("Word Documents", "*.docx"), 
                       ("Excel Files", "*.xlsx"), 
                       ("Text Files", "*.txt"),
                       ("PDF Files", "*.pdf"),
                       ("PowerPoint Files", "*.ppt;*.pptx")]
            # Prompt the user to select a file to upload
            file_path = filedialog.askopenfilename(filetypes=filetypes)
            if not file_path:
                return

            # Process the uploaded file and store its contents
            self.uploaded_file_data = self.process_file(file_path)

            # Display the file data in the text output widget
            self.text_output.insert(tk.END, "File data:\n")
            self.text_output.insert(tk.END, self.uploaded_file_data + "\n")
        except Exception as e:
            # Handle file upload errors
            messagebox.showerror("File Upload Error", f"An error occurred while uploading the file: {str(e)}")

    def process_file(self, file_path):
        """
        Process the uploaded file and return its content as a string.
        """
        try:
            # Determine file type based on file extension and process accordingly
            if file_path.endswith(".docx"):
                doc = Document(file_path)
                # Concatenate all paragraph text into a single string
                file_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            elif file_path.endswith(".xlsx"):
                df = pd.read_excel(file_path)
                # Convert the DataFrame to a string format
                file_content = df.to_string()
            elif file_path.endswith(".txt"):
                with open(file_path, "r") as file:
                    file_content = file.read()
            elif file_path.endswith(".pdf"):
                # Use PdfReader to read PDF files
                reader = PdfReader(file_path)
                file_content = ""
                for page in reader.pages:
                    file_content += page.extract_text()
            elif file_path.endswith(".pptx") or file_path.endswith(".ppt"):
                prs = Presentation(file_path)
                slides_content = []
                for slide in prs.slides:
                    slide_content = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_content += shape.text + "\n"
                    slides_content.append(slide_content)
                file_content = "\n".join(slides_content)
            else:
                file_content = "Unsupported file type."

            return file_content
        except Exception as e:
            # Handle file processing errors
            messagebox.showerror("File Processing Error", f"An error occurred while processing the file: {str(e)}")
            return ""

    def view_previous_conversations(self):
        """
        Display a new window showing previous conversations from the database.
        """
        # Create a new top-level window
        top = tk.Toplevel(self.root)
        top.title("Previous Conversations")

        # Create a frame to contain the conversation data
        frame = tk.Frame(top)
        frame.pack(padx=10, pady=10)

        # Create a label for the conversation data
        conversation_label = tk.Label(frame, text="Previous Conversations", font=("Arial", 12, "bold"))
        conversation_label.grid(row=0, column=0, columnspan=2, pady=5)

        # Create a scrollbar for the conversation data
        scrollbar = tk.Scrollbar(frame)
        scrollbar.grid(row=1, column=1, sticky="ns")

        # Create a listbox for displaying previous conversations
        listbox = tk.Listbox(frame, width=60, height=20, yscrollcommand=scrollbar.set)
        listbox.grid(row=1, column=0)

        # Link scrollbar to listbox
        scrollbar.config(command=listbox.yview)

        # Query the database to fetch previous conversations
        cursor = self.conn.execute("SELECT id, user_input, gpt_response FROM conversations")

        # Populate the listbox with previous conversations
        for row in cursor:
            conversation_id, user_input, gpt_response = row
            listbox.insert(tk.END, f"ID: {conversation_id} - User: {user_input} - ChatGPT: {gpt_response}")

        # Create a close button for the new window
        close_button = tk.Button(top, text="Close", command=top.destroy)
        close_button.pack(pady=5)

    def create_table_if_not_exists(self):
        """
        Create the conversations table in the database if it does not already exist.
        """
        with self.conn:
            # Create the conversations table
            self.conn.execute("""
                CREATE TABLE IF NOT EXISTS conversations (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_input TEXT,
                    gpt_response TEXT,
                    additional_info TEXT
                )
            """)

    def insert_conversation(self, user_input, gpt_response, additional_info):
        """
        Insert a conversation into the database.

        Args:
        - user_input (str): User input or prompt.
        - gpt_response (str): Generated response from ChatGPT.
        - additional_info (str): Additional information (optional).

        Returns:
        - bool: True if insertion is successful, False otherwise.
        """
        try:
            with self.conn:
                # Insert the conversation into the database
                self.conn.execute("INSERT INTO conversations (user_input, gpt_response, additional_info) VALUES (?, ?, ?)",
                                  (user_input, gpt_response, additional_info))
            return True
        except Exception as e:
            # Handle database insertion errors
            messagebox.showerror("Database Insertion Error", f"Failed to insert conversation into the database: {str(e)}")
            return False

    def close_app(self):
        """
        Close the application and database connection.
        """
        print("Exiting...")
        # Close the database connection
        self.conn.close()
        # Destroy the root window
        self.root.destroy()

def main():
    # Create a root window
    root = tk.Tk()
    # Initialize the ChatGPT application
    app = ChatGPTApp(root)
    # Start the main event loop
    root.mainloop()

# Run the program
if __name__ == "__main__":
    main()
